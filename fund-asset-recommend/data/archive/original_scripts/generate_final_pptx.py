from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

# 读取你最新上传的模板
template_path = '/Users/huangqingmeng/.openclaw/media/inbound/副本版本A_银行认股权业务培训PPT模版---7ad7bde6-3983-42e6-bdab-216df2f88bf.pptx'
template = Presentation(template_path)

# 创建新PPT，完全继承模板母版
prs = Presentation(template_path)

# 获取空白版式
blank_layout = None
for layout in template.slide_layouts:
    if len(layout.placeholders) == 0:
        blank_layout = layout
        break
if blank_layout is None:
    blank_layout = template.slide_layouts[6]

# 颜色定义
BLUE_DARK = RGBColor(0, 51, 102)
BLUE_MEDIUM = RGBColor(0, 102, 204)
BLACK = RGBColor(0, 0, 0)
RED = RGBColor(204, 0, 0)
LIGHT_BG = RGBColor(255, 255, 255)  # 内容页白色背景，匹配模板

# 分页处理函数
def add_content_page(section, current_title):
    lines = [l.rstrip() for l in section.split('\n') if l.rstrip()]
    slide = prs.slides.add_slide(blank_layout)
    
    # 设置白色背景，匹配模板
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    y = Inches(0.6)
    title_used = current_title
    
    for line in lines:
        line = line.rstrip()
        if not line.strip():
            continue
            
        if line.startswith('## '):
            text = line[3:].strip()
            title_used = text
            box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.5))
            tf = box.text_frame
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = BLUE_DARK
            y += Inches(0.7)
        elif line.startswith('### '):
            text = line[4:].strip()
            title_used = text
            box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.4))
            tf = box.text_frame
            p = tf.add_paragraph()
            p.text = text
            p.font.bold = True
            p.font.size = Pt(22)
            p.font.color.rgb = BLUE_MEDIUM
            y += Inches(0.5)
        elif line.startswith('|') and line.endswith('|'):
            # 表格行
            clean = re.sub(r'\|', '  ', line).strip()
            box = slide.shapes.add_textbox(Inches(1.0), y, Inches(11.0), Inches(0.28))
            tf = box.text_frame
            p = tf.add_paragraph()
            p.text = clean
            p.font.size = Pt(12)
            p.font.color.rgb = BLACK
            y += Inches(0.28)
            if y > Inches(6.5):
                y = Inches(0.6)
                slide = prs.slides.add_slide(blank_layout)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = LIGHT_BG
                if title_used:
                    box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.5))
                    tf = box.text_frame
